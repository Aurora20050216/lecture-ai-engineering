import os
import json
import numpy as np
import time
import random
from uuid import uuid4
from fastapi import FastAPI, Request, Body, Query
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pptx import Presentation
import boto3
import botocore.exceptions
from dotenv import load_dotenv
import traceback

# Load environment variables
load_dotenv()

# AWS settings
BUCKET_NAME = os.getenv("BUCKET_NAME", "your-default-bucket")
REGION = os.getenv("AWS_REGION", "us-east-1")

s3 = boto3.client("s3", region_name=REGION)
bedrock = boto3.client("bedrock-runtime", region_name=REGION)

# FastAPI setup
app = FastAPI()
templates = Jinja2Templates(directory="app/templates")
app.mount("/static", StaticFiles(directory="app/static"), name="static")

def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b) + 1e-10)

def get_embedding(text: str) -> np.ndarray:
    body = {
        "inputText": text
    }
    response = bedrock.invoke_model(
        modelId="amazon.titan-embed-text-v2:0",
        body=json.dumps(body),
        contentType="application/json",
        accept="application/json"
    )
    result = json.loads(response["body"].read())
    embedding = np.array(result.get("embedding", []))
    return embedding

def invoke_bedrock(slide_text: str, transcript: str) -> str:
    prompt = f"""
### Instruction:
以下は講義の文字起こしです。この文章には口語、言い直し、どもりなどが含まれています。
以下の制約のもとで、内容を要約せず、忠実に保ちながら、文法的に正しい自然な文章に整えてください。

- 話の順番、内容、意味を変えないでください。
- 話者の言い間違いや、口語表現は適切に直してください。
- 固有名詞の誤った文字起こしを修正する際はスライドを参考にしてください。
- 「えー」「あのー」などの無意味語は取り除いてください。
- 要約はせず、原文の流れを保ってください。
- タイムスタンプはそのままで良いです。
- 「以下は、講義の文字起こしを文法的に正しい自然な文章に整えたものです：」など言わずに、修正結果のみを返してください。
- これはある1つの長い文章を小分けにして送っています。最後に結合するので、毎回最後に改行を入れないようにしてもらえると、結合をしやすいです。

修正の例：
修正前：このようなモデルどのように学習しますかという話なんですが、ネクストtokenフリクションといって未来の位置一つ先のtoken単語を当てる学習をすることによって、言語モデルがうまく文章をモデル化する性能を高めていきます。
修正後：このようなモデルどのように学習しますかという話なんですが、"next token prediction"といって未来の一つ先のtoken(単語)を当てる学習をすることによって、言語モデルがうまく文章をモデル化する性能を高めていきます。

修正前：次に既存の既にアカデミアでよく使うられるようなベンチマークで性能がどれほど出てるかというところも確認します。
修正後：次に既存の既にアカデミアでよく使われるようなベンチマークで性能がどれほど出てるかというところも確認します。

修正前：さて前向きが前置きが長くなってしまったんですが、これから本題学習済みモデルを活用する技術についてお話します大きく分けて二つ話します。
修正後：さて前置きが長くなってしまったんですが、これから本題、学習済みモデルを活用する技術について大きく分けて二つお話しします。

### Slide Content:
{slide_text}

### Original Transcript:
{transcript}

### Output:
"""
    body = {
        "anthropic_version": "bedrock-2023-05-31",
         "messages": [
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 1024,
        "temperature": 0.2
    }
    response = bedrock.invoke_model(
        modelId="anthropic.claude-3-5-sonnet-20240620-v1:0",
        body=json.dumps(body),
        contentType="application/json",
        accept="application/json"
    )
    result = json.loads(response["body"].read())
    return "".join(part["text"] for part in result["content"] if part["type"] == "text").strip()

def split_text(text, max_len=500):
    return [text[i:i + max_len] for i in range(0, len(text), max_len)]


def invoke_bedrock_with_retry(slide_text, transcript, retries=20):
    for attempt in range(retries):
        try:
            return invoke_bedrock(slide_text, transcript)
        except botocore.exceptions.ClientError as e:
            if "Too many requests" in str(e):
                wait = (2 ** attempt) + random.uniform(0, 1)
                print(f"[Retry {attempt+1}] Waiting for {wait:.1f} seconds due to throttling...")
                time.sleep(wait)
            else:
                raise
    raise RuntimeError("Max retries reached. Failed to call Bedrock due to throttling.")

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


@app.get("/get_presigned_url")
async def get_presigned_url(filename: str = Query(...), content_type: str = Query(...)):
    try:
        presigned_url = s3.generate_presigned_url(
            ClientMethod='put_object',
            Params={'Bucket': BUCKET_NAME, 'Key': filename, 'ContentType': content_type},
            ExpiresIn=3600
        )
        return {"url": presigned_url}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/process_s3")
async def process_s3(pptx_key: str = Body(...), transcript_key: str = Body(...)):
    try:
        pptx_path = f"/tmp/{pptx_key.replace('/', '_')}"
        txt_path = f"/tmp/{transcript_key.replace('/', '_')}"

        s3.download_file(BUCKET_NAME, pptx_key, pptx_path)
        s3.download_file(BUCKET_NAME, transcript_key, txt_path)

        # Load transcript text
        with open(txt_path, "r", encoding="utf-8") as f:
            transcript = f.read()

        # Load slide texts
        prs = Presentation(pptx_path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides.append("\n".join(texts))

        # スライドテキストごとに埋め込み計算
        slide_embeddings = [get_embedding(s) for s in slides if s.strip()]

        # 文字起こしチャンク分割
        transcript_chunks = split_text(transcript, 500)

        # 文字起こしチャンクごとに埋め込み計算
        chunk_embeddings = [get_embedding(c) for c in transcript_chunks]

        # 各チャンクに対して最も類似度が高いスライドを対応付ける
        chunk_to_slide = []
        for chunk_emb in chunk_embeddings:
            sims = [cosine_similarity(chunk_emb, slide_emb) for slide_emb in slide_embeddings]
            best_slide_idx = int(np.argmax(sims))
            chunk_to_slide.append(best_slide_idx)

        # 対応付けたスライドテキストとチャンクのセットをBedrockで整形し、結果を結合
        fixed_chunks = []
        for idx, chunk in enumerate(transcript_chunks):
            slide_text = slides[chunk_to_slide[idx]]
            fixed = invoke_bedrock_with_retry(slide_text, chunk)
            fixed_chunks.append(fixed)

        cleaned_text = "\n\n".join(fixed_chunks)

        # Save output
        output_key = f"results/result_{uuid4()}.txt"
        output_path = f"/tmp/{output_key.replace('/', '_')}"

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(cleaned_text)

        s3.upload_file(output_path, BUCKET_NAME, output_key)

        presigned_download_url = s3.generate_presigned_url(
            ClientMethod='get_object',
            Params={'Bucket': BUCKET_NAME, 'Key': output_key},
            ExpiresIn=3600
        )

        return {"download_url": presigned_download_url}
    except Exception as e:
        print("エラー発生")
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)